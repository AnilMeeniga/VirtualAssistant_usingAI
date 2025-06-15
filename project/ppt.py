from pptx import Presentation
from pptx.util import Inches

# Create a PowerPoint presentation
prs = Presentation()

# Function to add a slide with a background
def add_slide_with_background(title_text, content_text, img_path=None):
    slide_layout = prs.slide_layouts[5]  # Title only layout
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = title_text

    # Add content
    textbox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(3))
    text_frame = textbox.text_frame
    text_frame.text = content_text

    # Add an image if provided
    if img_path:
        slide.shapes.add_picture(img_path, Inches(3), Inches(4), width=Inches(4))

# Add title slide with background
title_slide = prs.slides.add_slide(prs.slide_layouts[0])
title = title_slide.shapes.title
subtitle = title_slide.placeholders[1]
title.text = "Deep Learning Presentation"
subtitle.text = "An Introduction to Deep Learning"

# Add background image to title slide
bg_img_path = "background.jpg"  # Replace with your image path
title_slide.shapes.add_picture(bg_img_path, 0, 0, width=prs.slide_width, height=prs.slide_height)

# Overview slide
add_slide_with_background(
    "Overview of Deep Learning",
    "1. Introduction\n2. Key Concepts\n3. Neural Network Types\n4. Applications\n5. Challenges and Future",
    img_path="overview.png"  # Replace with your overview image
)

# Content slides with images
slides_content = [
    ("What is Deep Learning?", "Deep learning is a subset of machine learning that mimics the human brain.", "deep_learning.png"),
    ("History of Deep Learning", "Evolved from perceptrons in the 1950s to modern deep neural networks.", "history.png"),
    ("Key Concepts", "Neurons, Activation Functions, Layers, Weights, and Biases.", "concepts.png"),
    ("Types of Neural Networks", "CNN, RNN, GANs used for various tasks.", "network_types.png"),
    ("Deep Learning Frameworks", "Popular frameworks include TensorFlow, PyTorch, and Keras.", "frameworks.png"),
    ("Applications of Deep Learning", "Used in healthcare, finance, and autonomous vehicles.", "applications.png"),
    ("Challenges in Deep Learning", "Overfitting, data needs, interpretability, and cost.", "challenges.png"),
    ("Future of Deep Learning", "Explainable AI, better hardware, and edge AI.", "future.png"),
    ("Conclusion", "Deep learning continues to revolutionize industries.", None)
]

# Add slides with content and images
for slide in slides_content:
    add_slide_with_background(slide[0], slide[1], slide[2])

# Save the presentation
prs.save("Enhanced_Deep_Learning_Presentation.pptx")

print("Presentation with images and background created successfully!")